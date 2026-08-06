import os
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image, ImageDraw, ImageFont

OUTPUT_FILE = 'TS_OPAC_eLibrary_Feature_Report_v4.pptx'
IMAGE_DIR = 'presentation_images'

os.makedirs(IMAGE_DIR, exist_ok=True)

try:
    FONT_TITLE = ImageFont.truetype('arial.ttf', 24)
    FONT_NORMAL = ImageFont.truetype('arial.ttf', 16)
    FONT_SMALL = ImageFont.truetype('arial.ttf', 14)
except OSError:
    FONT_TITLE = ImageFont.load_default()
    FONT_NORMAL = ImageFont.load_default()
    FONT_SMALL = ImageFont.load_default()


def draw_card(draw, x, y, w, h, title, value=None, label=None, stroke='#CCCCCC'):
    draw.rectangle([x, y, x + w, y + h], fill='#FFFFFF', outline=stroke, width=2)
    draw.text((x + 16, y + 14), title, fill='#1F3859', font=FONT_NORMAL)
    if value is not None:
        draw.text((x + 16, y + 46), value, fill='#0B3D91', font=FONT_TITLE)
    if label:
        draw.text((x + 16, y + 90), label, fill='#5C6B73', font=FONT_SMALL)


def draw_table(draw, x, y, w, h, headers, rows):
    row_height = 38
    n_cols = len(headers)
    col_width = w // n_cols
    draw.rectangle([x, y, x + w, y + h], fill='#FFFFFF', outline='#A0A6AF', width=2)
    for col_idx, header in enumerate(headers):
        cell_x = x + col_idx * col_width
        draw.rectangle([cell_x, y, cell_x + col_width, y + row_height], fill='#F1F3F7', outline='#A0A6AF', width=1)
        draw.text((cell_x + 10, y + 10), header, fill='#1F3859', font=FONT_SMALL)
    for row_idx, row in enumerate(rows):
        row_y = y + (row_idx + 1) * row_height
        for col_idx, text in enumerate(row):
            cell_x = x + col_idx * col_width
            draw.text((cell_x + 10, row_y + 10), text, fill='#333333', font=FONT_SMALL)
            draw.rectangle([cell_x, row_y, cell_x + col_width, row_y + row_height], outline='#E1E5EA', width=1)


def draw_badge(draw, x, y, text, fill='#E7F1FF', text_fill='#0B3D91'):
    padding_x = 10
    padding_y = 6
    w = draw.textlength(text, font=FONT_SMALL) + padding_x * 2
    h = FONT_SMALL.size + padding_y * 2
    draw.rectangle([x, y, x + w, y + h], fill=fill, outline=None)
    draw.text((x + padding_x, y + padding_y), text, fill=text_fill, font=FONT_SMALL)
    return w, h


def draw_button(draw, x, y, w, h, text, fill='#0D6EFD', text_fill='#FFFFFF'):
    draw.rectangle([x, y, x + w, y + h], fill=fill, outline=None, width=0)
    draw.text((x + 14, y + (h - FONT_NORMAL.size) / 2), text, fill=text_fill, font=FONT_NORMAL)


def draw_table(draw, x, y, w, h, headers, rows, row_highlights=None):
    row_height = 40
    n_cols = len(headers)
    col_width = w // n_cols
    draw.rectangle([x, y, x + w, y + h], fill='#FFFFFF', outline='#BDC4CC', width=2)
    for col_idx, header in enumerate(headers):
        cell_x = x + col_idx * col_width
        draw.rectangle([cell_x, y, cell_x + col_width, y + row_height], fill='#F8FAFC', outline='#BDC4CC', width=1)
        draw.text((cell_x + 12, y + 12), header, fill='#1F3A72', font=FONT_SMALL)
    for row_idx, row in enumerate(rows):
        row_y = y + (row_idx + 1) * row_height
        for col_idx, text in enumerate(row):
            cell_x = x + col_idx * col_width
            if row_highlights and row_idx in row_highlights:
                draw.rectangle([cell_x, row_y, cell_x + col_width, row_y + row_height], fill=row_highlights[row_idx], outline='#BDC4CC', width=1)
            draw.text((cell_x + 12, row_y + 12), text, fill='#333333', font=FONT_SMALL)
            draw.rectangle([cell_x, row_y, cell_x + col_width, row_y + row_height], outline='#E2E8F0', width=1)


def create_staff_dashboard_image(filename):
    path = os.path.join(IMAGE_DIR, filename)
    img = Image.new('RGB', (1200, 720), '#F2F5F9')
    draw = ImageDraw.Draw(img)
    draw.text((36, 24), 'Staff Dashboard', fill='#0B3D91', font=FONT_TITLE)
    draw.text((36, 64), 'Circulation metrics and quick access to staff operations', fill='#5C6B73', font=FONT_SMALL)

    metrics = [
        ('Active Loans', '128', '#0D6EFD'),
        ('Overdue Items', '19', '#DC3545'),
        ('Holds Waiting', '42', '#0DCAF0'),
        ('Ready for Pickup', '11', '#198754'),
        ('In Transit', '7', '#6C757D'),
        ('Pending Requests', '5', '#FFC107'),
    ]
    for idx, (title, value, color) in enumerate(metrics):
        x = 36 + (idx % 3) * 380
        y = 110 + (idx // 3) * 130
        draw.rectangle([x, y, x + 340, y + 110], fill='#FFFFFF', outline='#CED4DA', width=2)
        draw.text((x + 16, y + 14), title, fill='#495057', font=FONT_NORMAL)
        draw.text((x + 16, y + 44), value, fill=color, font=FONT_TITLE)
        draw.text((x + 16, y + 82), 'Metric', fill='#6C757D', font=FONT_SMALL)

    draw.rectangle([36, 360, 1164, 450], fill='#FFFFFF', outline='#CED4DA', width=2)
    draw.text((50, 374), 'Quick Access', fill='#0B3D91', font=FONT_NORMAL)
    buttons = [
        ('Circulation Hub', '#0D6EFD'),
        ('Check Out', '#198754'),
        ('Check In', '#0DC263'),
        ('Borrowers', '#FFC107'),
        ('Holds', '#6C757D'),
    ]
    for idx, (label, color) in enumerate(buttons):
        draw_button(draw, 46 + idx * 224, 404, 210, 42, label, fill=color)

    draw.text((36, 470), 'Recent Checkouts (Last 24 Hours)', fill='#0B3D91', font=FONT_NORMAL)
    checkouts = [
        ('Alice Roberts', 'The Hobbit', 'Aug 12'),
        ('Mark Liu', 'Python Basics', 'Aug 15'),
        ('Sofia Lee', 'Library Science', 'Aug 17'),
    ]
    draw_table(draw, 36, 500, 560, 140, ['Borrower', 'Item', 'Due Date'], checkouts)

    draw.text((616, 470), 'Recent Returns (Last 24 Hours)', fill='#0B3D91', font=FONT_NORMAL)
    returns = [
        ('Jenna Park', 'Django Unchained', 'Aug 07'),
        ('Omar Hassan', 'Data Structures', 'Aug 07'),
        ('Nina Das', 'Modern Catalogs', 'Aug 06'),
    ]
    draw_table(draw, 616, 500, 560, 140, ['Borrower', 'Item', 'Returned'], returns)

    draw.rectangle([36, 660, 560, 700], fill='#FFFFFF', outline='#CED4DA', width=2)
    draw.text((50, 672), 'Trend chart placeholder', fill='#6C757D', font=FONT_SMALL)
    draw.rectangle([620, 660, 1144, 700], fill='#FFFFFF', outline='#CED4DA', width=2)
    draw.text((634, 672), 'Loan status distribution placeholder', fill='#6C757D', font=FONT_SMALL)

    img.save(path)
    return path


def create_admin_dashboard_image(filename):
    path = os.path.join(IMAGE_DIR, filename)
    img = Image.new('RGB', (1200, 720), '#F8F9FA')
    draw = ImageDraw.Draw(img)
    draw.text((36, 24), 'Administrator Dashboard', fill='#1F3A72', font=FONT_TITLE)
    draw.text((36, 64), 'System overview, status tables, collection metrics, and admin controls', fill='#5C6B73', font=FONT_SMALL)

    cards = [
        ('Total Users', '432', '#0D6EFD'),
        ('Active Borrowers', '287', '#0DCAF0'),
        ('Publications', '5,124', '#198754'),
        ('Total Items', '12,598', '#FFC107'),
    ]
    for idx, (title, value, color) in enumerate(cards):
        x = 36 + idx * 285
        draw.rectangle([x, 110, x + 260, 190], fill='#FFFFFF', outline=color, width=2)
        draw.text((x + 16, 124), title, fill='#495057', font=FONT_NORMAL)
        draw.text((x + 16, 158), value, fill=color, font=FONT_TITLE)

    draw.text((36, 220), 'System Status', fill='#1F3A72', font=FONT_NORMAL)
    status_rows = [
        ('Active Loans', '128'),
        ('Overdue Items', '19'),
        ('Blocked Borrowers', '6'),
        ('Pending Holds', '42'),
    ]
    draw_table(draw, 36, 250, 560, 160, ['Metric', 'Value'], status_rows)

    draw.text((616, 220), 'Collection Status', fill='#1F3A72', font=FONT_NORMAL)
    collection_rows = [
        ('Available Items', '8,762'),
        ('Items on Loan', '2,184'),
        ('Items on Hold Shelf', '342'),
        ('Items in Transit', '7'),
    ]
    draw_table(draw, 616, 250, 560, 160, ['Metric', 'Value'], collection_rows)

    control_panels = [
        ('User Management', 'Manage users and roles', '#DC3545'),
        ('System Configuration', 'Configure policies and lookup', '#FFC107'),
        ('System Reports', 'View detailed reports', '#0DCAF0'),
    ]
    for idx, (title, desc, color) in enumerate(control_panels):
        x = 36 + idx * 385
        y = 430
        draw.rectangle([x, y, x + 360, y + 160], fill='#FFFFFF', outline='#CED4DA', width=2)
        draw.text((x + 16, y + 16), title, fill=color, font=FONT_NORMAL)
        draw.text((x + 16, y + 52), desc, fill='#5C6B73', font=FONT_SMALL)
        draw_button(draw, x + 16, y + 100, 180, 36, 'Open', fill=color)

    draw.rectangle([36, 610, 1164, 690], fill='#FFFFFF', outline='#CED4DA', width=2)
    draw.text((50, 628), 'Administrator panel note: Manage users, settings, and system reports from this dashboard.', fill='#495057', font=FONT_SMALL)

    img.save(path)
    return path


def create_borrower_account_image(filename):
    path = os.path.join(IMAGE_DIR, filename)
    img = Image.new('RGB', (1200, 720), '#F7F8FA')
    draw = ImageDraw.Draw(img)
    draw.text((36, 24), 'My Account', fill='#205081', font=FONT_TITLE)
    draw.text((36, 64), 'Borrower self-service view with current loans, holds, requests, and history', fill='#5C6B73', font=FONT_SMALL)
    draw.text((36, 92), 'Welcome, John Doe!', fill='#495057', font=FONT_NORMAL)

    draw.rectangle([36, 130, 420, 300], fill='#FFFFFF', outline='#CED4DA', width=2)
    draw.text((56, 146), 'Account Information', fill='#1F3A72', font=FONT_NORMAL)
    account_items = [
        'Username: john_doe',
        'Email: john@example.com',
        'Card Number: 1234567890',
        'Active Loans: 3 / 5',
    ]
    for idx, text in enumerate(account_items):
        draw.text((56, 176 + idx * 24), text, fill='#333333', font=FONT_NORMAL)
    draw_button(draw, 56, 270, 140, 32, 'Edit Profile', fill='#0D6EFD')

    draw.text((456, 134), 'Current Loans', fill='#205081', font=FONT_NORMAL)
    loans = [
        ('The Hobbit', 'Aug 12', 'Renew'),
        ('Python Basics', 'Aug 15', 'Renew'),
        ('Library Science', 'Aug 17', 'Cannot renew'),
    ]
    draw_table(draw, 456, 160, 712, 140, ['Title', 'Due Date', 'Actions'], loans, row_highlights={2: '#FFF3CD'})
    draw_badge(draw, 600, 180, 'Overdue 2 day(s)', fill='#F8D7DA', text_fill='#842029')

    draw.text((36, 320), 'Active Holds', fill='#205081', font=FONT_NORMAL)
    holds = [
        ('Django Unchained', 'Ready for Pickup', 'Main Hall Shelf'),
        ('Modern Catalogs', 'Waiting', 'Queue #3'),
    ]
    draw_table(draw, 36, 350, 560, 100, ['Title', 'Status', 'Details'], holds)
    draw_text_y = 350 + 100 + 16
    draw.text((36, draw_text_y), 'Action: Cancel hold', fill='#DC3545', font=FONT_SMALL)

    draw.text((36, 470), 'Checkout Requests', fill='#205081', font=FONT_NORMAL)
    requests = [
        ('Data Structures', 'Pending Review', 'Requested Aug 05'),
        ('Library Metrics', 'Approved', 'Pickup by Aug 10'),
    ]
    draw_table(draw, 36, 500, 560, 120, ['Title', 'Status', 'Details'], requests)
    draw_text_y = 500 + 120 + 16
    draw.text((36, draw_text_y), 'Approved request shows pickup location and date', fill='#6C757D', font=FONT_SMALL)

    draw.text((456, 320), 'Loan History', fill='#205081', font=FONT_NORMAL)
    history = [
        ('Advanced Python', '07/01/2026', '07/15/2026'),
        ('Library Science', '06/10/2026', '06/24/2026'),
    ]
    draw_table(draw, 456, 350, 712, 120, ['Title', 'Checkout Date', 'Return Date'], history)

    img.save(path)
    return path


def create_search_image(filename):
    path = os.path.join(IMAGE_DIR, filename)
    img = Image.new('RGB', (1200, 700), '#F2F5F9')
    draw = ImageDraw.Draw(img)
    draw.text((40, 24), 'OPAC Search and Publication Record', fill='#1F3A72', font=FONT_TITLE)
    draw.text((40, 64), 'Search input, filters, publication title, author, location, and availability', fill='#5C6B73', font=FONT_SMALL)

    draw.rectangle([40, 120, 1160, 170], fill='#FFFFFF', outline='#C0CAD3', width=2)
    draw.text((60, 136), 'Search keyword: title, author, subject, ISBN...', fill='#6A737D', font=FONT_NORMAL)

    draw.rectangle([40, 190, 1160, 245], fill='#FFFFFF', outline='#C0CAD3', width=2)
    filter_label = 'Filter: Language | Type | Location | Availability'
    draw.text((60, 206), filter_label, fill='#0B3D91', font=FONT_NORMAL)

    draw.rectangle([40, 270, 1160, 620], fill='#FFFFFF', outline='#C0CAD3', width=2)
    draw.text((60, 292), 'Publication: Python Basics', fill='#1F3A72', font=FONT_NORMAL)
    draw.text((60, 322), 'Author: Jane Developer', fill='#333333', font=FONT_NORMAL)
    draw.text((60, 352), 'Call Number: QA76.73.P98', fill='#333333', font=FONT_NORMAL)
    draw.text((60, 382), 'Location: 2nd Floor - Computer Science', fill='#333333', font=FONT_NORMAL)
    draw.text((60, 412), 'Availability: Available', fill='#198754', font=FONT_NORMAL)
    draw.text((60, 442), 'Summary: Practical introduction to Python programming.', fill='#5C6B73', font=FONT_SMALL)

    img.save(path)
    return path


prs = Presentation()


def add_notes(slide, notes):
    slide.notes_slide.notes_text_frame.text = notes

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = 'TS OPAC eLibrary Feature Presentation'
subtitle = slide.placeholders[1]
subtitle.text = 'Actual dashboard and borrower account mockups with system labels'
add_notes(slide, 'Introduce the deck with visuals that reflect the real page elements and labels from the application templates.')

# Staff Dashboard slide
staff_image = create_staff_dashboard_image('staff_dashboard_v5.png')
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'Staff Dashboard'
slide.shapes.placeholders[1].text = 'Active loans, overdue items, holds, quick actions, and recent activity'
slide.shapes.add_picture(staff_image, Inches(0.5), Inches(1.8), width=Inches(9))
add_notes(slide, 'Show the staff dashboard with exact metric names and quick access labels from the real system.')

# Admin Dashboard slide
admin_image = create_admin_dashboard_image('admin_dashboard_v5.png')
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'Administrator Dashboard'
slide.shapes.placeholders[1].text = 'System overview cards, status tables, collection metrics, and admin controls'
slide.shapes.add_picture(admin_image, Inches(0.5), Inches(1.8), width=Inches(9))
add_notes(slide, 'Show the administrator dashboard with system and collection status details similar to the live templates.')

# Borrower Account slide
borrower_image = create_borrower_account_image('borrower_account_v5.png')
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'Borrower Account Page'
slide.shapes.placeholders[1].text = 'Account summary, current loans, active holds, and checkout requests'
slide.shapes.add_picture(borrower_image, Inches(0.5), Inches(1.8), width=Inches(9))
add_notes(slide, 'Show borrower self-service functionality with loan, hold, and request details.')

# Search and Publication slide
search_image = create_search_image('search_record_v5.png')
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'OPAC Search and Publication Detail'
slide.shapes.placeholders[1].text = 'Search box, filter bar, and publication metadata view'
slide.shapes.add_picture(search_image, Inches(0.5), Inches(1.8), width=Inches(9))
add_notes(slide, 'Show the OPAC search input and publication record metadata that correspond to the actual catalog template.')

# Feature Highlights slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'Feature Highlights'
body = slide.shapes.placeholders[1].text_frame
body.text = 'Implemented features'
for point in [
    'Staff dashboard with circulation metrics and quick action buttons',
    'Administrator dashboard with system and collection summaries',
    'Borrower account page with current loans, holds, and checkout request tracking',
    'OPAC search with keyword and filter support, plus publication availability',
    'Notifications for due soon, overdue, and hold ready scenarios',
]:
    p = body.add_paragraph()
    p.text = point
    p.level = 0
    p.font.size = Pt(18)
add_notes(slide, 'Summarize the main system features using the actual page labels.')

# Next Steps slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'Next Steps'
body = slide.shapes.placeholders[1].text_frame
body.text = 'Recommended follow-up work'
for point in [
    'Capture real browser screenshots for staff and borrower pages',
    'Validate the login/CAPTCHA workflow in the live app',
    'Enable scheduled notifications and email delivery',
    'Document demo account credentials and feature mapping',
]:
    p = body.add_paragraph()
    p.text = point
    p.level = 0
    p.font.size = Pt(18)
add_notes(slide, 'Close with concrete next actions aligned to the actual application views.')

prs.save(OUTPUT_FILE)
print(f'Created {OUTPUT_FILE}')
